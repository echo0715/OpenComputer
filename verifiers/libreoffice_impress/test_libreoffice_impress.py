"""
Test LibreOffice Impress verifier endpoints in a live E2B sandbox.

Covers:
  - Help/usage output
  - Error cases (LibreOffice not running, bad args)
  - ODF file parsing endpoints (ODP file, no UNO needed)
  - UNO live endpoints with PPTX file (slides, shapes, text, transitions, notes)
  - Composite check-* endpoints — positive and negative cases

Usage:
    python verifiers/libreoffice_impress/test_libreoffice_impress.py
"""

import json
import sys
import time
import traceback
from pathlib import Path

from dotenv import load_dotenv
from e2b_desktop import Sandbox
from e2b.sandbox.commands.command_handle import CommandExitException

load_dotenv()

VERIFIER_LOCAL = Path(__file__).parent / "libreoffice_impress.py"
VERIFIER_REMOTE = "/home/user/verifiers/libreoffice_impress.py"
V = f"python3 {VERIFIER_REMOTE}"

TEST_ODP_PATH = "/home/user/test_verifier.odp"
TEST_PPTX_PATH = "/home/user/test_verifier.pptx"

# ── ODP fixture (stdlib ZIP+XML) ────────────────────────────────────────

CREATE_ODP_SCRIPT = r'''
import zipfile

CONTENT_XML = """<?xml version="1.0" encoding="UTF-8"?>
<office:document-content
  xmlns:office="urn:oasis:names:tc:opendocument:xmlns:office:1.0"
  xmlns:text="urn:oasis:names:tc:opendocument:xmlns:text:1.0"
  xmlns:draw="urn:oasis:names:tc:opendocument:xmlns:drawing:1.0"
  xmlns:presentation="urn:oasis:names:tc:opendocument:xmlns:presentation:1.0"
  xmlns:svg="urn:oasis:names:tc:opendocument:xmlns:svg-compatible:1.0"
  office:version="1.2">
  <office:body>
    <office:presentation>
      <draw:page draw:name="TitleSlide" draw:master-page-name="Default">
        <draw:frame draw:name="title" presentation:class="title" svg:width="20cm" svg:height="3cm" svg:x="2cm" svg:y="5cm">
          <draw:text-box>
            <text:p>Project Overview</text:p>
          </draw:text-box>
        </draw:frame>
        <draw:frame draw:name="subtitle" presentation:class="subtitle" svg:width="20cm" svg:height="5cm" svg:x="2cm" svg:y="10cm">
          <draw:text-box>
            <text:p>A comprehensive look at our progress</text:p>
          </draw:text-box>
        </draw:frame>
      </draw:page>
      <draw:page draw:name="ContentSlide" draw:master-page-name="Default">
        <draw:frame draw:name="title" presentation:class="title" svg:width="20cm" svg:height="3cm" svg:x="2cm" svg:y="1cm">
          <draw:text-box>
            <text:p>Key Findings</text:p>
          </draw:text-box>
        </draw:frame>
        <draw:frame draw:name="content" presentation:class="object" svg:width="20cm" svg:height="12cm" svg:x="2cm" svg:y="5cm">
          <draw:text-box>
            <text:p>Revenue grew by 25%</text:p>
            <text:p>Customer satisfaction improved</text:p>
          </draw:text-box>
        </draw:frame>
      </draw:page>
      <draw:page draw:name="ConclusionSlide" draw:master-page-name="Default">
        <draw:frame draw:name="title" svg:width="20cm" svg:height="3cm" svg:x="2cm" svg:y="5cm">
          <draw:text-box>
            <text:p>Conclusion</text:p>
          </draw:text-box>
        </draw:frame>
      </draw:page>
    </office:presentation>
  </office:body>
</office:document-content>
"""

MANIFEST_XML = """<?xml version="1.0" encoding="UTF-8"?>
<manifest:manifest xmlns:manifest="urn:oasis:names:tc:opendocument:xmlns:manifest:1.0">
  <manifest:file-entry manifest:media-type="application/vnd.oasis.opendocument.presentation" manifest:full-path="/"/>
  <manifest:file-entry manifest:media-type="text/xml" manifest:full-path="content.xml"/>
</manifest:manifest>
"""

path = "/home/user/test_verifier.odp"
with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as z:
    z.writestr("mimetype", "application/vnd.oasis.opendocument.presentation")
    z.writestr("content.xml", CONTENT_XML)
    z.writestr("META-INF/manifest.xml", MANIFEST_XML)

print("OK")
'''

# ── PPTX fixture (python-pptx) ──────────────────────────────────────────

CREATE_PPTX_SCRIPT = r'''
import subprocess, sys
subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx", "-q"],
                      stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

prs = Presentation()

# Slide 1: Title slide
slide1 = prs.slides.add_slide(prs.slide_layouts[0])
slide1.shapes.title.text = "Annual Report 2025"
slide1.placeholders[1].text = "Company Overview and Results"

# Slide 2: Content slide
slide2 = prs.slides.add_slide(prs.slide_layouts[1])
slide2.shapes.title.text = "Revenue Summary"
body = slide2.placeholders[1]
tf = body.text_frame
tf.text = "Q1: $1.2M"
tf.add_paragraph().text = "Q2: $1.5M"
tf.add_paragraph().text = "Q3: $1.8M"
tf.add_paragraph().text = "Q4: $2.1M"

# Slide 3: Blank slide with shapes
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
# Add a text box
txBox = slide3.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(2))
tf = txBox.text_frame
tf.text = "Custom Text Box Content"

# Slide 4: Another content slide
slide4 = prs.slides.add_slide(prs.slide_layouts[1])
slide4.shapes.title.text = "Conclusion"
slide4.placeholders[1].text = "Strong performance across all metrics"

# Add speaker notes to slide 1
notes_slide = slide1.notes_slide
notes_slide.notes_text_frame.text = "Welcome everyone to the annual report presentation"

prs.save("/home/user/test_verifier.pptx")
print("OK")
'''

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

passed = 0
failed = 0
errors: list[str] = []


class CmdResult:
    def __init__(self, exit_code: int, stdout: str, stderr: str):
        self.exit_code = exit_code
        self.stdout = stdout
        self.stderr = stderr


def run(sandbox: Sandbox, cmd: str, timeout: int = 30) -> dict | list:
    r = run_raw(sandbox, cmd, timeout)
    if r.exit_code != 0 and not r.stdout.strip():
        return {"error": f"exit_code={r.exit_code} stderr={r.stderr[:300]}"}
    try:
        return json.loads(r.stdout)
    except json.JSONDecodeError:
        return {"error": f"Invalid JSON: {r.stdout[:300]}"}


def run_raw(sandbox: Sandbox, cmd: str, timeout: int = 30) -> CmdResult:
    try:
        result = sandbox.commands.run(f"{V} {cmd}", timeout=timeout)
        return CmdResult(result.exit_code, result.stdout, result.stderr)
    except CommandExitException as e:
        return CmdResult(e.exit_code, e.stdout, e.stderr)


def run_shell(sandbox: Sandbox, cmd: str, timeout: int = 60) -> CmdResult:
    try:
        result = sandbox.commands.run(cmd, timeout=timeout)
        return CmdResult(result.exit_code, result.stdout, result.stderr)
    except CommandExitException as e:
        return CmdResult(e.exit_code, e.stdout, e.stderr)


def check(name: str, condition: bool, detail: str = ""):
    global passed, failed
    if condition:
        passed += 1
        print(f"  PASS  {name}")
    else:
        failed += 1
        msg = f"  FAIL  {name}"
        if detail:
            msg += f"  -- {detail}"
        print(msg)
        errors.append(f"{name}: {detail}")


def is_valid_json(stdout: str) -> bool:
    try:
        json.loads(stdout)
        return True
    except (json.JSONDecodeError, TypeError):
        return False


# ---------------------------------------------------------------------------
# Test groups
# ---------------------------------------------------------------------------

def test_help(sandbox: Sandbox):
    print("\n=== Help ===")
    result = run_raw(sandbox, "--help")
    check("help exits 0", result.exit_code == 0, f"got exit_code={result.exit_code}")
    check("help mentions commands", "Commands:" in result.stdout, result.stdout[:100])


def test_errors_lo_not_running(sandbox: Sandbox):
    print("\n=== Errors (LibreOffice not running) ===")
    for cmd, name in [("slides", "slides"), ("slide-text 0", "slide-text"), ("doc-info", "doc-info")]:
        data = run(sandbox, cmd)
        check(f"{name} returns error when LO down", "error" in data, str(data)[:100])


def test_errors_bad_args(sandbox: Sandbox):
    print("\n=== Errors (bad args) ===")
    result = run_raw(sandbox, "check-slide-title")
    check("missing arg exits 1", result.exit_code == 1)
    check("missing arg valid JSON", is_valid_json(result.stdout), result.stdout[:100])

    result = run_raw(sandbox, "nonexistent-command")
    check("unknown cmd exits 1", result.exit_code == 1)
    check("unknown cmd valid JSON", is_valid_json(result.stdout), result.stdout[:100])


def test_odp_file_parsing(sandbox: Sandbox):
    """ODP parsing endpoints (no UNO needed)."""
    print("\n=== ODP File Parsing ===")

    data = run(sandbox, f"parse-slides {TEST_ODP_PATH}")
    check("parse-slides returns slides", len(data.get("slides", [])) > 0, str(data)[:150])
    if "error" not in data:
        check("parse-slides has 3 slides", data.get("count") == 3, f"count={data.get('count')}")
        names = [s["name"] for s in data.get("slides", [])]
        check("parse-slides has TitleSlide", "TitleSlide" in names, str(names))

    data = run(sandbox, f"parse-slide-text 0 {TEST_ODP_PATH}")
    check("parse-slide-text returns texts", "texts" in data and "error" not in data, str(data)[:150])
    if "error" not in data:
        check("parse-slide-text has Project Overview",
              "Project Overview" in data.get("full_text", ""),
              data.get("full_text", "")[:100])

    data = run(sandbox, f"parse-slide-text 1 {TEST_ODP_PATH}")
    if "error" not in data:
        check("parse-slide-text slide 2 has Key Findings",
              "Key Findings" in data.get("full_text", ""),
              data.get("full_text", "")[:100])

    data = run(sandbox, "parse-slides /nonexistent/file.odp")
    check("parse missing file returns error", "error" in data, str(data)[:100])


def test_uno_basic(sandbox: Sandbox):
    """Basic UNO endpoints with PPTX loaded."""
    print("\n=== UNO Basic (PPTX) ===")

    data = run(sandbox, "slides")
    check("slides returns dict", isinstance(data, dict) and "error" not in data, str(data)[:200])
    if "error" not in data:
        check("slides count = 4", data.get("count") == 4, f"count={data.get('count')}")

    data = run(sandbox, "doc-info")
    check("doc-info has title", "title" in data and "error" not in data, str(data)[:150])
    if "error" not in data:
        check("doc-info slide_count = 4", data.get("slide_count") == 4,
              f"slide_count={data.get('slide_count')}")

    data = run(sandbox, "slide-size")
    check("slide-size returns dict", isinstance(data, dict) and "error" not in data, str(data)[:100])
    if "error" not in data:
        check("slide-size has width", "width" in data, str(data.keys()))

    data = run(sandbox, "master-slides")
    check("master-slides returns dict", isinstance(data, dict) and "error" not in data, str(data)[:100])


def test_uno_slide_content(sandbox: Sandbox):
    """Slide text and shape endpoints via UNO."""
    print("\n=== UNO Slide Content ===")

    # Slide 1 text (title slide)
    data = run(sandbox, "slide-text 0")
    check("slide-text 0 returns texts", "texts" in data and "error" not in data, str(data)[:200])
    if "error" not in data:
        full = data.get("full_text", "")
        check("slide 0 has Annual Report", "Annual Report" in full, full[:100])

    # Slide 2 text (content slide)
    data = run(sandbox, "slide-text 1")
    if "error" not in data:
        full = data.get("full_text", "")
        check("slide 1 has Revenue", "Revenue" in full, full[:100])

    # Shapes on slide 1
    data = run(sandbox, "slide-shapes 0")
    check("slide-shapes returns shapes", "shapes" in data and "error" not in data, str(data)[:200])
    if "error" not in data:
        check("slide 0 has shapes", data.get("count", 0) > 0, f"count={data.get('count')}")

    # Layout
    data = run(sandbox, "slide-layout 0")
    check("slide-layout returns dict", isinstance(data, dict) and "error" not in data, str(data)[:150])
    if "error" not in data:
        check("slide-layout has layout", "layout" in data, str(data.keys()))
        check("slide-layout has master_page", "master_page" in data, str(data.keys()))


def test_uno_notes_and_transitions(sandbox: Sandbox):
    """Notes and transition endpoints."""
    print("\n=== UNO Notes & Transitions ===")

    # Notes on slide 1
    data = run(sandbox, "notes 0")
    check("notes returns dict", isinstance(data, dict) and "error" not in data, str(data)[:200])
    if "error" not in data:
        check("notes has content", "notes" in data, str(data.keys()))
        # Our PPTX has notes on slide 1
        check("notes slide 0 has text", len(data.get("notes", "")) > 0,
              f"notes={data.get('notes', '')[:80]}")

    # Transition
    data = run(sandbox, "transition 0")
    check("transition returns dict", isinstance(data, dict) and "error" not in data, str(data)[:100])
    if "error" not in data:
        check("transition has type", "type" in data, str(data.keys()))


def test_checks_positive(sandbox: Sandbox):
    """Composite check-* endpoints — positive cases."""
    print("\n=== Checks (positive) ===")

    # check-slide-count
    data = run(sandbox, "check-slide-count 4")
    check("check-slide-count 4 match", data.get("match") is True, str(data)[:100])

    # check-slide-title
    data = run(sandbox, 'check-slide-title 0 "Annual Report"')
    check("check-slide-title match", data.get("match") is True, str(data)[:150])

    # check-slide-contains
    data = run(sandbox, 'check-slide-contains 1 "Q1"')
    check("check-slide-contains Q1", data.get("contains") is True, str(data)[:100])

    # check-shape-count (slide 0 has at least some shapes)
    shapes = run(sandbox, "slide-shapes 0")
    if "error" not in shapes:
        count = shapes["count"]
        data = run(sandbox, f"check-shape-count 0 {count}")
        check("check-shape-count match", data.get("match") is True, str(data)[:100])

    # check-file-exists
    data = run(sandbox, f"check-file-exists {TEST_PPTX_PATH}")
    check("check-file-exists true", data.get("exists") is True, str(data)[:100])

    # check-file-saved
    data = run(sandbox, "check-file-saved")
    check("check-file-saved has key", "saved" in data or "error" in data, str(data)[:100])


def test_checks_negative(sandbox: Sandbox):
    """Composite check-* endpoints — negative cases."""
    print("\n=== Checks (negative) ===")

    data = run(sandbox, "check-slide-count 99")
    check("check-slide-count mismatch", data.get("match") is False, str(data)[:100])

    data = run(sandbox, 'check-slide-title 0 "NonexistentTitle12345"')
    check("check-slide-title mismatch", data.get("match") is False, str(data)[:100])

    data = run(sandbox, 'check-slide-contains 0 "xyzzy_nonexistent_12345"')
    check("check-slide-contains false", data.get("contains") is False, str(data)[:100])

    data = run(sandbox, "check-file-exists /nonexistent/file.pptx")
    check("check-file-exists false", data.get("exists") is False, str(data)[:100])

    data = run(sandbox, "check-shape-count 0 9999")
    check("check-shape-count mismatch", data.get("match") is False, str(data)[:100])


def test_all_commands_return_json(sandbox: Sandbox):
    """Every CLI command should output valid JSON."""
    print("\n=== JSON validity (all commands) ===")

    no_arg_cmds = ["slides", "doc-info", "slide-size", "master-slides"]
    for cmd in no_arg_cmds:
        result = run_raw(sandbox, cmd)
        valid = is_valid_json(result.stdout)
        check(f"{cmd} valid JSON", valid,
              f"exit={result.exit_code} stdout={result.stdout[:80]}" if not valid else "")

    arg_cmds = [
        ("slide-text", "0"),
        ("slide-text", "1"),
        ("slide-shapes", "0"),
        ("slide-layout", "0"),
        ("notes", "0"),
        ("transition", "0"),
        ("parse-slides", TEST_ODP_PATH),
        ("parse-slide-text", f"0 {TEST_ODP_PATH}"),
        ("check-slide-count", "4"),
        ("check-slide-title", '0 "test"'),
        ("check-slide-contains", '0 "test"'),
        ("check-shape-count", "0 1"),
        ("check-has-transition", "0"),
        ("check-file-exists", TEST_PPTX_PATH),
        ("check-file-saved", ""),
    ]

    for cmd, arg in arg_cmds:
        full_cmd = f"{cmd} {arg}".strip()
        result = run_raw(sandbox, full_cmd)
        valid = is_valid_json(result.stdout)
        check(f"{full_cmd} valid JSON", valid,
              f"exit={result.exit_code} stdout={result.stdout[:80]}" if not valid else "")


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def launch_libreoffice(sandbox: Sandbox, file_path: str) -> bool:
    """Launch LibreOffice Impress with UNO socket, opening the given file."""
    try:
        sandbox.commands.run("pkill -9 -f soffice", timeout=5)
    except (CommandExitException, Exception):
        pass
    try:
        sandbox.commands.run("pkill -9 -f oosplash", timeout=5)
    except (CommandExitException, Exception):
        pass
    time.sleep(2)

    launch_script = f'''#!/bin/bash
soffice --impress \
  '--accept=socket,host=localhost,port=2002;urp;' \
  --norestore --nologo \
  {file_path} \
  > /tmp/lo.log 2>&1
'''
    sandbox.files.write("/tmp/launch_lo.sh", launch_script)
    sandbox.commands.run("chmod +x /tmp/launch_lo.sh", timeout=3)
    sandbox.commands.run("/tmp/launch_lo.sh", background=True)

    check_script = '''
import uno, time
for attempt in range(5):
    try:
        ctx = uno.getComponentContext()
        resolver = ctx.ServiceManager.createInstanceWithContext(
            "com.sun.star.bridge.UnoUrlResolver", ctx
        )
        resolved = resolver.resolve(
            "uno:socket,host=localhost,port=2002;urp;StarOffice.ComponentContext"
        )
        smgr = resolved.ServiceManager
        desktop = smgr.createInstanceWithContext("com.sun.star.frame.Desktop", resolved)
        doc = desktop.getCurrentComponent()
        if doc is not None:
            print("READY")
            break
        else:
            print("NODOC")
            break
    except Exception as e:
        if attempt < 4:
            time.sleep(1)
        else:
            print(f"ERR:{e}")
'''
    sandbox.files.write("/tmp/check_uno.py", check_script)

    for i in range(30):
        try:
            r = sandbox.commands.run("python3 /tmp/check_uno.py", timeout=10)
            if "READY" in r.stdout:
                print(f"  LibreOffice UNO ready (attempt {i+1})")
                return True
            if "NODOC" in r.stdout:
                print(f"  UNO connected but no doc yet (attempt {i+1})")
        except (CommandExitException, Exception) as e:
            if i > 15:
                print(f"  attempt {i+1}: {str(e)[:80]}")
        time.sleep(2)

    try:
        r = sandbox.commands.run("tail -20 /tmp/lo.log", timeout=3)
        print(f"  [lo log] {r.stdout[-300:]}")
    except (CommandExitException, Exception):
        pass
    return False


def main():
    global passed, failed

    print("=" * 60)
    print("LibreOffice Impress Verifier Test Suite")
    print("=" * 60)

    print("\nCreating sandbox from desktop-all-apps...")
    sandbox = Sandbox.create(template="desktop-all-apps", timeout=600)

    try:
        # Upload verifier
        print(f"Uploading {VERIFIER_LOCAL} -> {VERIFIER_REMOTE}")
        sandbox.commands.run("mkdir -p /home/user/verifiers")
        with open(VERIFIER_LOCAL) as f:
            sandbox.files.write(VERIFIER_REMOTE, f.read())

        # ── Tests with LibreOffice NOT running ──
        test_help(sandbox)
        test_errors_lo_not_running(sandbox)
        test_errors_bad_args(sandbox)

        # ── Create ODP test file ──
        print("\n  Creating ODP test file...")
        sandbox.files.write("/home/user/create_odp.py", CREATE_ODP_SCRIPT)
        r = run_shell(sandbox, "python3 /home/user/create_odp.py")
        odp_ok = "OK" in r.stdout
        check("ODP file created", odp_ok, f"stdout={r.stdout[:100]} stderr={r.stderr[:100]}")
        if odp_ok:
            test_odp_file_parsing(sandbox)

        # ── Create PPTX test file ──
        print("\n  Creating PPTX test file...")
        sandbox.files.write("/home/user/create_pptx.py", CREATE_PPTX_SCRIPT)
        r = run_shell(sandbox, "python3 /home/user/create_pptx.py", timeout=120)
        pptx_ok = "OK" in r.stdout
        check("PPTX file created", pptx_ok, f"stdout={r.stdout[:100]} stderr={r.stderr[:100]}")

        # ── Launch LibreOffice Impress with PPTX ──
        print("\nLaunching LibreOffice Impress with PPTX file...")
        lo_ready = launch_libreoffice(sandbox, TEST_PPTX_PATH)
        if not lo_ready:
            print("  WARNING: LibreOffice UNO not ready -- UNO tests may fail")

        # ── UNO tests ──
        test_uno_basic(sandbox)
        test_uno_slide_content(sandbox)
        test_uno_notes_and_transitions(sandbox)
        test_checks_positive(sandbox)
        test_checks_negative(sandbox)
        test_all_commands_return_json(sandbox)

    except Exception:
        traceback.print_exc()
        failed += 1
        errors.append(f"Unhandled exception: {traceback.format_exc()}")

    finally:
        sandbox.kill()
        print("\nSandbox killed.")

    # ── Summary ──
    print("\n" + "=" * 60)
    print(f"Results: {passed} passed, {failed} failed, {passed + failed} total")
    print("=" * 60)

    if errors:
        print("\nFailures:")
        for e in errors:
            print(f"  - {e}")

    sys.exit(1 if failed > 0 else 0)


if __name__ == "__main__":
    main()
